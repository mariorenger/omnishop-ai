"""File ingestion: turn an uploaded file (many types) into plain text for RAG.

Text extraction is best-effort and degrades gracefully. Heavy/optional libs are
imported lazily so a missing one only disables that one file type. Scanned PDFs
and images route through the swappable OCRProvider.
"""
from __future__ import annotations
import io
import json
import os

TEXT_EXT = {".txt", ".md", ".markdown", ".log", ".rst", ".text"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff", ".gif"}
SUPPORTED = TEXT_EXT | IMAGE_EXT | {".csv", ".tsv", ".json", ".html", ".htm",
                                   ".pdf", ".docx", ".pptx", ".xlsx"}


def _media_type(ext: str) -> str:
    return {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".webp": "image/webp", ".gif": "image/gif", ".bmp": "image/bmp",
            ".tif": "image/tiff", ".tiff": "image/tiff"}.get(ext, "image/png")


def _decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except Exception:  # noqa: BLE001
            continue
    return data.decode("utf-8", errors="ignore")


def extract_text(filename: str, data: bytes, ocr=None) -> str:
    ext = os.path.splitext(filename or "")[1].lower()

    if ext in TEXT_EXT or ext in (".csv", ".tsv"):
        return _decode(data).strip()

    if ext == ".json":
        try:
            return json.dumps(json.loads(_decode(data)), ensure_ascii=False, indent=2)
        except Exception:  # noqa: BLE001
            return _decode(data).strip()

    if ext in (".html", ".htm"):
        try:
            from bs4 import BeautifulSoup
            return BeautifulSoup(_decode(data), "html.parser").get_text("\n").strip()
        except Exception:
            return _decode(data).strip()

    if ext == ".docx":
        from docx import Document
        doc = Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts).strip()

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
        return "\n".join(parts).strip()

    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts).strip()

    if ext == ".pdf":
        text = _pdf_text(data)
        if len(text.strip()) >= 20 or ocr is None:
            return text.strip()
        # likely scanned -> OCR each page image if poppler/pdf2image present
        return _pdf_ocr(data, ocr).strip() or text.strip()

    if ext in IMAGE_EXT:
        if ocr is None:
            return ""
        return ocr.extract_image(data, _media_type(ext)).strip()

    # unknown: try to decode as text, else reject
    text = _decode(data).strip()
    if text:
        return text
    raise ValueError(f"unsupported file type: {ext or '(none)'}")


def _pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    except Exception:  # noqa: BLE001
        return ""


def _pdf_ocr(data: bytes, ocr) -> str:
    try:
        from pdf2image import convert_from_bytes
        images = convert_from_bytes(data, dpi=200)
    except Exception as e:  # noqa: BLE001 — poppler/pdf2image missing
        return f"[scanned PDF; OCR rasterization unavailable: {e}]"
    out = []
    for img in images:
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        out.append(ocr.extract_image(buf.getvalue(), "image/png"))
    return "\n".join(out)
